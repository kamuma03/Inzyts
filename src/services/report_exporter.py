"""
Report Exporter Service.

Orchestrates report generation in multiple formats (HTML, PDF, PPTX, Markdown)
from completed analysis notebooks. Integrates PII detection and executive
summary generation into the export pipeline.
"""

import base64
import io
import re
import time
from datetime import datetime
from enum import Enum
from pathlib import Path
from typing import Any, Dict, List, Optional

import markdown as md_lib
import nbformat
from jinja2 import Environment, FileSystemLoader, select_autoescape
from pydantic import BaseModel

from src.config import settings
from src.services.executive_summary import ExecutiveSummary, ExecutiveSummaryGenerator
from src.services.pii_detector import PIIDetector, PIIScanResult
from src.utils.path_validator import ensure_dir
from src.utils.logger import get_logger

logger = get_logger()


class ReportFormat(str, Enum):
    HTML = "html"
    PDF = "pdf"
    PPTX = "pptx"
    MARKDOWN = "markdown"


class ReportResult(BaseModel):
    """Metadata about a generated report."""

    format: ReportFormat
    file_path: str
    file_size_bytes: int
    generation_time_seconds: float
    has_executive_summary: bool
    has_pii_warnings: bool


class _CellContent(BaseModel):
    """Processed cell content for template rendering."""

    cell_type: str
    source: str = ""
    html_content: str = ""  # Rendered markdown → HTML
    text_output: str = ""
    images: List[str] = []  # Base64-encoded PNG images


class _Section(BaseModel):
    """A logical section of the report."""

    title: str = ""
    cells: List[_CellContent] = []


# Inzyts brand colors for PPTX
_TURQUOISE = (45, 212, 191)
_COBALT = (30, 58, 95)
_COBALT_LIGHT = (42, 79, 122)
_WHITE = (255, 255, 255)
_TEXT_SECONDARY = (100, 116, 139)


class ReportExporter:
    """Generates reports in HTML, PDF, PPTX, and Markdown formats."""

    def __init__(self):
        self._summary_generator = ExecutiveSummaryGenerator()
        self._pii_detector = PIIDetector()
        self._jinja_env = Environment(
            loader=FileSystemLoader(Path(__file__).parent / "templates"),
            autoescape=select_autoescape(["html"]),
        )

    def export(
        self,
        notebook_path: str,
        job_metadata: Dict[str, Any],
        format: ReportFormat = ReportFormat.HTML,
        include_executive_summary: bool = True,
        include_pii_masking: bool = False,
    ) -> ReportResult:
        """Generate a report from a completed analysis notebook.

        Args:
            notebook_path: Path to the .ipynb file.
            job_metadata: Dict with keys: job_id, mode, title, question, model, version.
            format: Target export format.
            include_executive_summary: Whether to generate executive summary.
            include_pii_masking: Whether to mask detected PII in output.

        Returns:
            ReportResult with file path and metadata.
        """
        start = time.monotonic()

        # 1. Read notebook
        nb = self._read_notebook(notebook_path)
        if nb is None:
            raise ValueError(f"Could not read notebook at {notebook_path}")

        # 2. Run PII scan
        pii_result = self._pii_detector.scan_notebook(notebook_path)

        # 3. Generate executive summary
        exec_summary: Optional[ExecutiveSummary] = None
        if include_executive_summary:
            try:
                exec_summary = self._summary_generator.generate(notebook_path)
            except Exception as e:
                logger.warning(f"Executive summary generation failed: {e}")

        # 4. Process notebook cells into sections
        sections = self._parse_sections(nb, include_pii_masking)

        # 5. Extract execution metrics from the last section
        metrics = self._extract_metrics(nb)

        # 6. Build output path
        report_dir = Path(settings.output_dir) / "reports"
        ensure_dir(report_dir)
        job_id = job_metadata.get("job_id", "unknown")
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        ext = "md" if format == ReportFormat.MARKDOWN else format.value
        output_path = report_dir / f"report_{job_id}_{timestamp}.{ext}"

        # 7. Render to target format
        if format == ReportFormat.HTML:
            self._render_html(
                output_path, sections, metrics, exec_summary,
                pii_result, job_metadata,
            )
        elif format == ReportFormat.PDF:
            self._render_pdf(
                output_path, sections, metrics, exec_summary,
                pii_result, job_metadata,
            )
        elif format == ReportFormat.PPTX:
            self._render_pptx(
                output_path, nb, exec_summary, pii_result,
                metrics, job_metadata,
            )
        elif format == ReportFormat.MARKDOWN:
            self._render_markdown(
                output_path, nb, exec_summary, pii_result,
                metrics, job_metadata, include_pii_masking,
            )

        elapsed = time.monotonic() - start
        file_size = output_path.stat().st_size

        return ReportResult(
            format=format,
            file_path=str(output_path),
            file_size_bytes=file_size,
            generation_time_seconds=round(elapsed, 2),
            has_executive_summary=exec_summary is not None,
            has_pii_warnings=pii_result.has_pii,
        )

    # -------------------------------------------------------------------------
    # Notebook parsing
    # -------------------------------------------------------------------------

    def _read_notebook(self, path: str) -> Optional[nbformat.NotebookNode]:
        try:
            with open(path, "r", encoding="utf-8") as f:
                return nbformat.read(f, as_version=4)
        except Exception as e:
            logger.error(f"Failed to read notebook: {e}")
            return None

    def _parse_sections(
        self, nb: nbformat.NotebookNode, mask_pii: bool = False
    ) -> List[_Section]:
        """Parse notebook cells into logical sections based on ## headers."""
        sections: List[_Section] = []
        current_section = _Section()

        for cell in nb.cells:
            source = cell.source
            if mask_pii:
                source = self._pii_detector.mask_text(source)

            # Check for section headers (## N. Title)
            header_match = re.search(r"^##\s+\d+\.\s+(.+)", source, re.MULTILINE)
            if cell.cell_type == "markdown" and header_match:
                # Save previous section if it has content
                if current_section.cells or current_section.title:
                    sections.append(current_section)
                current_section = _Section(title=header_match.group(1).strip())
                # Add remaining content after header
                remaining = source[header_match.end():].strip()
                if remaining:
                    current_section.cells.append(
                        _CellContent(
                            cell_type="markdown",
                            source=remaining,
                            html_content=md_lib.markdown(
                                remaining, extensions=["tables", "fenced_code"]
                            ),
                        )
                    )
                continue

            processed = self._process_cell(cell, source, mask_pii)
            current_section.cells.append(processed)

        # Don't forget the last section
        if current_section.cells or current_section.title:
            sections.append(current_section)

        return sections

    def _process_cell(
        self, cell: nbformat.NotebookNode, source: str, mask_pii: bool
    ) -> _CellContent:
        """Process a single notebook cell into renderable content."""
        if cell.cell_type == "markdown":
            html = md_lib.markdown(source, extensions=["tables", "fenced_code"])
            return _CellContent(
                cell_type="markdown", source=source, html_content=html
            )

        # Code cell
        text_output = ""
        images: List[str] = []

        if hasattr(cell, "outputs"):
            for output in cell.outputs:
                otype = output.get("output_type", "")
                if otype == "stream":
                    text = output.get("text", "")
                    if mask_pii:
                        text = self._pii_detector.mask_text(text)
                    text_output += text
                elif otype in ("execute_result", "display_data"):
                    data = output.get("data", {})
                    if "image/png" in data:
                        images.append(data["image/png"])
                    if "text/plain" in data:
                        text = data["text/plain"]
                        if mask_pii:
                            text = self._pii_detector.mask_text(text)
                        text_output += text
                elif otype == "error":
                    text_output += "\n".join(output.get("traceback", []))

        return _CellContent(
            cell_type="code",
            source=source,
            text_output=text_output.strip(),
            images=images,
        )

    def _extract_metrics(self, nb: nbformat.NotebookNode) -> Dict[str, str]:
        """Extract execution metrics table from the last markdown cell."""
        metrics: Dict[str, str] = {}
        for cell in reversed(nb.cells):
            if cell.cell_type != "markdown":
                continue
            if "Execution Metrics" not in cell.source:
                continue
            # Parse markdown table
            for line in cell.source.split("\n"):
                line = line.strip()
                if line.startswith("|") and "---" not in line:
                    parts = [p.strip() for p in line.split("|") if p.strip()]
                    if len(parts) == 2 and parts[0] != "Metric":
                        metrics[parts[0]] = parts[1]
            break
        return metrics

    # -------------------------------------------------------------------------
    # HTML rendering
    # -------------------------------------------------------------------------

    def _render_html(
        self,
        output_path: Path,
        sections: List[_Section],
        metrics: Dict[str, str],
        executive_summary: Optional[ExecutiveSummary],
        pii_result: Optional[PIIScanResult],
        job_metadata: Dict[str, Any],
    ) -> None:
        template = self._jinja_env.get_template("report.html")
        html = template.render(
            title=job_metadata.get("title", "Analysis Report"),
            mode=job_metadata.get("mode", "Exploratory"),
            generated_at=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            model=job_metadata.get("model"),
            version=job_metadata.get("version"),
            question=job_metadata.get("question"),
            executive_summary=executive_summary,
            pii_result=pii_result,
            sections=[s.model_dump() for s in sections],
            metrics=metrics,
        )
        output_path.write_text(html, encoding="utf-8")

    # -------------------------------------------------------------------------
    # PDF rendering (via WeasyPrint)
    # -------------------------------------------------------------------------

    def _render_pdf(
        self,
        output_path: Path,
        sections: List[_Section],
        metrics: Dict[str, str],
        executive_summary: Optional[ExecutiveSummary],
        pii_result: Optional[PIIScanResult],
        job_metadata: Dict[str, Any],
    ) -> None:
        try:
            from weasyprint import HTML as WPHTML
        except ImportError:
            raise RuntimeError(
                "PDF export requires weasyprint. "
                "Install with: pip install weasyprint"
            )

        # Render HTML first, then convert to PDF
        template = self._jinja_env.get_template("report.html")
        html = template.render(
            title=job_metadata.get("title", "Analysis Report"),
            mode=job_metadata.get("mode", "Exploratory"),
            generated_at=datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
            model=job_metadata.get("model"),
            version=job_metadata.get("version"),
            question=job_metadata.get("question"),
            executive_summary=executive_summary,
            pii_result=pii_result,
            sections=[s.model_dump() for s in sections],
            metrics=metrics,
        )
        WPHTML(string=html).write_pdf(str(output_path))

    # -------------------------------------------------------------------------
    # PPTX rendering (via python-pptx)
    # -------------------------------------------------------------------------

    def _render_pptx(
        self,
        output_path: Path,
        nb: nbformat.NotebookNode,
        executive_summary: Optional[ExecutiveSummary],
        pii_result: Optional[PIIScanResult],
        metrics: Dict[str, str],
        job_metadata: Dict[str, Any],
    ) -> None:
        try:
            from pptx import Presentation
            from pptx.dml.color import RGBColor
            from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
            from pptx.util import Inches, Pt, Emu
        except ImportError:
            raise RuntimeError(
                "PPTX export requires python-pptx. "
                "Install with: pip install python-pptx"
            )

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        # Helper functions
        def add_blank_slide():
            layout = prs.slide_layouts[6]  # Blank layout
            return prs.slides.add_slide(layout)

        def add_textbox(slide, left, top, width, height):
            return slide.shapes.add_textbox(left, top, width, height)

        def set_paragraph(para, text, size=14, bold=False, color=_TEXT_SECONDARY, alignment=PP_ALIGN.LEFT):
            para.text = text
            para.font.size = Pt(size)
            para.font.bold = bold
            para.font.color.rgb = RGBColor(*color)
            para.alignment = alignment

        def add_slide_header(slide, title_text):
            """Add a consistent header bar to a slide."""
            # Header background
            shape = slide.shapes.add_shape(
                1, Inches(0), Inches(0), slide_width, Inches(1.1)
            )
            shape.fill.solid()
            shape.fill.fore_color.rgb = RGBColor(*_COBALT)
            shape.line.fill.background()
            # Title text
            txbox = add_textbox(slide, Inches(0.6), Inches(0.2), Inches(10), Inches(0.7))
            p = txbox.text_frame.paragraphs[0]
            set_paragraph(p, title_text, size=28, bold=True, color=_WHITE)
            # Brand accent line
            line = slide.shapes.add_shape(
                1, Inches(0.6), Inches(1.0), Inches(2), Inches(0.06)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor(*_TURQUOISE)
            line.line.fill.background()

        def add_bullet_list(slide, items, left, top, width, height, size=14, color=_TEXT_SECONDARY):
            """Add a bullet-point list to a slide."""
            txbox = add_textbox(slide, left, top, width, height)
            tf = txbox.text_frame
            tf.word_wrap = True
            for i, item in enumerate(items):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = item
                p.font.size = Pt(size)
                p.font.color.rgb = RGBColor(*color)
                p.space_after = Pt(6)
                p.level = 0

        # ---- Slide 1: Title ----
        slide = add_blank_slide()
        # Full-slide gradient background
        bg_shape = slide.shapes.add_shape(
            1, Inches(0), Inches(0), slide_width, slide_height
        )
        bg_shape.fill.solid()
        bg_shape.fill.fore_color.rgb = RGBColor(*_COBALT)
        bg_shape.line.fill.background()

        # Brand name
        txbox = add_textbox(slide, Inches(1), Inches(1.5), Inches(10), Inches(0.8))
        p = txbox.text_frame.paragraphs[0]
        set_paragraph(p, "INZYTS", size=20, bold=True, color=_TURQUOISE)

        # Title
        txbox = add_textbox(slide, Inches(1), Inches(2.3), Inches(10), Inches(1.5))
        p = txbox.text_frame.paragraphs[0]
        set_paragraph(
            p,
            job_metadata.get("title", "Analysis Report"),
            size=40, bold=True, color=_WHITE,
        )

        # Metadata
        meta_lines = [
            f"Mode: {job_metadata.get('mode', 'Exploratory')}",
            f"Date: {datetime.now().strftime('%Y-%m-%d')}",
        ]
        if job_metadata.get("model"):
            meta_lines.append(f"Model: {job_metadata['model']}")
        if job_metadata.get("question"):
            meta_lines.append(f"Question: {job_metadata['question']}")

        txbox = add_textbox(slide, Inches(1), Inches(4.2), Inches(10), Inches(2))
        tf = txbox.text_frame
        for i, line_text in enumerate(meta_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            set_paragraph(p, line_text, size=16, color=_WHITE)

        # Accent line
        accent = slide.shapes.add_shape(
            1, Inches(1), Inches(4.0), Inches(3), Inches(0.06)
        )
        accent.fill.solid()
        accent.fill.fore_color.rgb = RGBColor(*_TURQUOISE)
        accent.line.fill.background()

        # ---- Slide 2: Executive Summary ----
        if executive_summary:
            slide = add_blank_slide()
            add_slide_header(slide, "Executive Summary")

            # Summary text
            if executive_summary.summary_text:
                txbox = add_textbox(slide, Inches(0.6), Inches(1.3), Inches(12), Inches(1.5))
                tf = txbox.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                # Truncate if too long for a slide
                text = executive_summary.summary_text
                if len(text) > 600:
                    text = text[:597] + "..."
                set_paragraph(p, text, size=14, color=_TEXT_SECONDARY)

            # Key findings
            if executive_summary.key_findings:
                txbox = add_textbox(slide, Inches(0.6), Inches(3.0), Inches(5.5), Inches(0.5))
                p = txbox.text_frame.paragraphs[0]
                set_paragraph(p, "Key Findings", size=18, bold=True, color=_COBALT)
                add_bullet_list(
                    slide, executive_summary.key_findings,
                    Inches(0.6), Inches(3.5), Inches(5.5), Inches(3.5),
                )

            # Recommendations
            if executive_summary.recommendations:
                txbox = add_textbox(slide, Inches(7), Inches(3.0), Inches(5.5), Inches(0.5))
                p = txbox.text_frame.paragraphs[0]
                set_paragraph(p, "Recommendations", size=18, bold=True, color=_COBALT)
                add_bullet_list(
                    slide, executive_summary.recommendations,
                    Inches(7), Inches(3.5), Inches(5.5), Inches(3.5),
                )

        # ---- Slide 3: Data Quality ----
        if executive_summary and executive_summary.data_quality_highlights:
            slide = add_blank_slide()
            add_slide_header(slide, "Data Quality Overview")
            add_bullet_list(
                slide, executive_summary.data_quality_highlights,
                Inches(0.6), Inches(1.5), Inches(11), Inches(4),
                size=16,
            )

        # ---- Slides 4+: Notebook sections with charts ----
        current_section_title = ""
        section_images: List[str] = []
        section_text_parts: List[str] = []

        for cell in nb.cells:
            # Check for section header
            header_match = re.match(
                r"^##\s+\d+\.\s+(.+)", cell.source, re.MULTILINE
            )
            if cell.cell_type == "markdown" and header_match:
                # Flush previous section
                if current_section_title and (section_text_parts or section_images):
                    self._add_section_slide(
                        prs, add_blank_slide, add_slide_header,
                        add_bullet_list, add_textbox,
                        current_section_title, section_text_parts,
                        section_images, slide_width,
                    )
                current_section_title = header_match.group(1).strip()
                section_text_parts = []
                section_images = []

                # Add remaining text from header cell
                remaining = cell.source[header_match.end():].strip()
                if remaining:
                    for line in remaining.split("\n"):
                        line = line.strip()
                        if line and not line.startswith("---") and not line.startswith("|"):
                            clean = re.sub(r"^[-*#]+\s*", "", line)
                            if clean and len(clean) > 5:
                                section_text_parts.append(clean)
                continue

            if cell.cell_type == "markdown":
                for line in cell.source.split("\n"):
                    line = line.strip()
                    if line and not line.startswith("---") and not line.startswith("|"):
                        clean = re.sub(r"^[-*#]+\s*", "", line)
                        if clean and len(clean) > 5:
                            section_text_parts.append(clean)

            elif cell.cell_type == "code" and hasattr(cell, "outputs"):
                for output in cell.outputs:
                    if output.get("output_type") in ("display_data", "execute_result"):
                        data = output.get("data", {})
                        if "image/png" in data:
                            section_images.append(data["image/png"])

        # Flush last section
        if current_section_title and (section_text_parts or section_images):
            self._add_section_slide(
                prs, add_blank_slide, add_slide_header,
                add_bullet_list, add_textbox,
                current_section_title, section_text_parts,
                section_images, slide_width,
            )

        # ---- PII Warning Slide ----
        if pii_result and pii_result.has_pii:
            slide = add_blank_slide()
            add_slide_header(slide, "PII Detection Warning")
            findings_text = [
                f"{f.pii_type.upper()} ({f.severity}) - {f.location}"
                for f in pii_result.findings[:10]
            ]
            if len(pii_result.findings) > 10:
                findings_text.append(
                    f"... and {len(pii_result.findings) - 10} more"
                )
            add_bullet_list(
                slide, findings_text,
                Inches(0.6), Inches(1.5), Inches(11), Inches(5),
            )

        # ---- Last Slide: Metrics ----
        if metrics:
            slide = add_blank_slide()
            add_slide_header(slide, "Execution Metrics")

            from pptx.util import Inches as PInches

            # Build a simple table
            rows = len(metrics) + 1
            cols = 2
            table_shape = slide.shapes.add_table(
                rows, cols,
                Inches(1), Inches(1.5), Inches(6), Inches(0.4 * rows + 0.4),
            )
            table = table_shape.table

            # Header row
            for i, header in enumerate(["Metric", "Value"]):
                cell = table.cell(0, i)
                cell.text = header
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(12)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(*_WHITE)
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(*_COBALT)

            # Data rows
            for row_idx, (key, value) in enumerate(metrics.items(), start=1):
                for col_idx, text in enumerate([key, value]):
                    cell = table.cell(row_idx, col_idx)
                    cell.text = text
                    for p in cell.text_frame.paragraphs:
                        p.font.size = Pt(11)
                        p.font.color.rgb = RGBColor(*_TEXT_SECONDARY)

        prs.save(str(output_path))

    def _add_section_slide(
        self, prs, add_blank_slide, add_slide_header,
        add_bullet_list, add_textbox,
        title: str, text_parts: List[str],
        images: List[str], slide_width,
    ) -> None:
        """Add a content slide for a notebook section."""
        from pptx.util import Inches

        slide = add_blank_slide()
        add_slide_header(slide, title)

        # Limit bullet points to prevent overflow
        bullets = text_parts[:8]
        if len(text_parts) > 8:
            bullets.append(f"... and {len(text_parts) - 8} more points")

        has_image = len(images) > 0

        if has_image:
            # Split layout: text on left, image on right
            add_bullet_list(
                slide, bullets,
                Inches(0.6), Inches(1.3), Inches(5.5), Inches(5.5),
                size=13,
            )
            # Add first chart image
            try:
                img_data = base64.b64decode(images[0])
                img_stream = io.BytesIO(img_data)
                slide.shapes.add_picture(
                    img_stream, Inches(6.5), Inches(1.3), Inches(6.3), Inches(5.5)
                )
            except Exception as e:
                logger.warning(f"Failed to embed chart image in PPTX: {e}")
        else:
            add_bullet_list(
                slide, bullets,
                Inches(0.6), Inches(1.3), Inches(12), Inches(5.5),
            )

        # Add extra slides for additional images
        for img_b64 in images[1:4]:  # Max 3 extra image slides
            try:
                extra_slide = add_blank_slide()
                add_slide_header(extra_slide, f"{title} (continued)")
                img_data = base64.b64decode(img_b64)
                img_stream = io.BytesIO(img_data)
                extra_slide.shapes.add_picture(
                    img_stream, Inches(1.5), Inches(1.5), Inches(10), Inches(5.5)
                )
            except Exception as e:
                logger.warning(f"Failed to embed additional chart in PPTX: {e}")

    # -------------------------------------------------------------------------
    # Markdown rendering
    # -------------------------------------------------------------------------

    def _render_markdown(
        self,
        output_path: Path,
        nb: nbformat.NotebookNode,
        executive_summary: Optional[ExecutiveSummary],
        pii_result: Optional[PIIScanResult],
        metrics: Dict[str, str],
        job_metadata: Dict[str, Any],
        mask_pii: bool = False,
    ) -> None:
        parts: List[str] = []

        # Header
        title = job_metadata.get("title", "Analysis Report")
        parts.append(f"# {title}")
        parts.append("")
        parts.append(f"**Mode:** {job_metadata.get('mode', 'Exploratory')} | "
                      f"**Date:** {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}")
        if job_metadata.get("model"):
            parts.append(f"**Model:** {job_metadata['model']}")
        if job_metadata.get("question"):
            parts.append(f"**Question:** {job_metadata['question']}")
        parts.append("")

        # PII Warning
        if pii_result and pii_result.has_pii:
            parts.append("---")
            parts.append("")
            parts.append(f"> **PII Warning:** {len(pii_result.findings)} potential "
                          f"PII item(s) detected.")
            parts.append("")

        # Executive Summary
        if executive_summary:
            parts.append("---")
            parts.append("")
            parts.append("## Executive Summary")
            parts.append("")
            if executive_summary.summary_text:
                parts.append(executive_summary.summary_text)
                parts.append("")
            parts.append("### Key Findings")
            for f in executive_summary.key_findings:
                parts.append(f"- {f}")
            parts.append("")
            parts.append("### Recommendations")
            for r in executive_summary.recommendations:
                parts.append(f"- {r}")
            parts.append("")
            if executive_summary.data_quality_highlights:
                parts.append("### Data Quality")
                for h in executive_summary.data_quality_highlights:
                    parts.append(f"- {h}")
                parts.append("")

        # Notebook content
        parts.append("---")
        parts.append("")
        for cell in nb.cells:
            source = cell.source
            if mask_pii:
                source = self._pii_detector.mask_text(source)

            if cell.cell_type == "markdown":
                parts.append(source)
                parts.append("")
            elif cell.cell_type == "code":
                parts.append("```python")
                parts.append(source)
                parts.append("```")
                parts.append("")
                # Include text outputs
                if hasattr(cell, "outputs"):
                    for output in cell.outputs:
                        otype = output.get("output_type", "")
                        if otype == "stream":
                            text = output.get("text", "")
                            if mask_pii:
                                text = self._pii_detector.mask_text(text)
                            if text.strip():
                                parts.append("```")
                                parts.append(text.strip())
                                parts.append("```")
                                parts.append("")

        # Metrics
        if metrics:
            parts.append("---")
            parts.append("")
            parts.append("## Execution Metrics")
            parts.append("")
            parts.append("| Metric | Value |")
            parts.append("|--------|-------|")
            for key, value in metrics.items():
                parts.append(f"| {key} | {value} |")
            parts.append("")

        # Footer
        parts.append("---")
        parts.append(f"*Generated by Inzyts — {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}*")

        output_path.write_text("\n".join(parts), encoding="utf-8")
